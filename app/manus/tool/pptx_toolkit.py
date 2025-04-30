# Copyright 2025 ZTE Corporation.
# All Rights Reserved.
#
#    Licensed under the Apache License, Version 2.0 (the "License"); you may
#    not use this file except in compliance with the License. You may obtain
#    a copy of the License at
#
#         http://www.apache.org/licenses/LICENSE-2.0
#
#    Unless required by applicable law or agreed to in writing, software
#    distributed under the License is distributed on an "AS IS" BASIS, WITHOUT
#    WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied. See the
#    License for the specific language governing permissions and limitations
#    under the License.

from pptx import Presentation


def extract_pptx_content(pptx_path: str) -> str:
    print(f'read {pptx_path}')
    prs = Presentation(pptx_path)

    all_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_content.append(shape.text)
            elif shape.has_text_frame:
                all_content.append(shape.text_frame.text)
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    all_content.append(" | ".join(row_text))

    return '\n'.join(all_content)


if __name__ == '__main__':
    result = extract_pptx_content("/home/eval_test.pptx")
    print(result)
